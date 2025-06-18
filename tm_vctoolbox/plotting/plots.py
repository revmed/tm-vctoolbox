"""
File containing general purpose plotting code
"""

# %%
import matplotlib.pyplot as plt

# Sample data for the table
data = [["A", 1, 10], ["B", 2, 20], ["C", 3, 30]]
columns = ["Column 1", "Column 2", "Column 3"]

fig, ax = plt.subplots(figsize=(8, 2))  # Adjust figsize as needed
ax.axis("off")  # Hide axes for a cleaner table appearance
table = ax.table(cellText=data, colLabels=columns, loc="center")
table.auto_set_font_size(False)
table.set_fontsize(10)
table.scale(1.2, 1.2)  # Adjust scaling if necessary

plt.savefig(
    "my_table.png", bbox_inches="tight", dpi=300
)  # Save with high DPI for better quality

# %%
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide_layout = prs.slide_layouts[
    5
]  # Choose a suitable slide layout (e.g., Title and Content)
slide = prs.slides.add_slide(slide_layout)

left = top = width = height = Inches(1)  # Adjust position and size as needed
pic = slide.shapes.add_picture(
    "my_table.png", left, top, width=Inches(6)
)  # Adjust width/height

prs.save("presentation_with_table.pptx")

# %%
import matplotlib.pyplot as plt

# Set the fonttype parameter for searchable text in PDF
plt.rcParams["pdf.fonttype"] = 42
plt.rcParams["ps.fonttype"] = 42  # Also set for PostScript if needed

# Create your plot
fig, ax = plt.subplots()
ax.plot([0, 1, 2], [0, 1, 4])
ax.set_title("Searchable Text Example")
ax.set_xlabel("X-axis Label")
ax.set_ylabel("Y-axis Label")
ax.text(0.5, 2, "Some text on the plot", horizontalalignment="center")

# Save the figure as a PDF
plt.savefig("searchable_figure.pdf", bbox_inches="tight")

plt.show()

# %%
